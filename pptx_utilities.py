from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Pt
from common.survey_structure import SurveyStructure
#rom data_collector import DataCollector

levels = { 1: "Непосредственный руководитель",
           2: "Вышестоящий руководитель",
           3: "Функциональный руководитель"}

levels_table = { 1: "Прямые подчиненные",
                 2: "Подчиненные ниже уровнем",
                 3: "Функциональные подчиненные"}

levels_chart = { 1: "Прямое подчинение",
                 2: "Нижнеуровневое подчинение",
                 3: "Функциональное подчинение"}

levels_rate_column = { 1: 4,
                       2: 11,
                       3: 20}

def update_cards(slide, report, level : str):
    survey_structure = SurveyStructure("resources\\survey_structure.json")
    selects = [i for i in survey_structure["groups"][level] if survey_structure["fields"][i][0] == "select"]
    headers = [survey_structure["output headers"][i] for i in selects]

    positive_counter = 0
    negative_counter = 0
    for header in headers:
        if report[header + "_count"] == 0:
            continue

        if report[header + "_positive_pct"] >= 0.7:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text == "p" + str(positive_counter):
                    update_text(shape.text_frame, header)  
                    positive_counter += 1
                    break
        else:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text == "n" + str(negative_counter):
                    update_text(shape.text_frame, header)
                    negative_counter += 1
                    break

    for i in range(positive_counter, 5):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text == "p" + str(i):
                slide.shapes.element.remove(shape.element)  
    for i in range(negative_counter, 5):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text == "n" + str(i):
                slide.shapes.element.remove(shape.element)  

def update_text(text_frame, text):
    text_frame.paragraphs[0].runs[0].text = text
    if text == "Недостаточно данных":
        text_frame.paragraphs[0].runs[0].font.size = Pt(26)

    for i, run in enumerate(text_frame.paragraphs[0].runs,0):
        if i > 0:
            run.text = ""

def update_table_cell(table, row : int, column : int, value : str):
    update_text(table.cell(row, column).text_frame, value)

def update_chart(chart, new_data):
    series0 = list()
    series1 = list()
    for pair in new_data[1]:
        series0.append(pair[0])
        series1.append(pair[1])
    
    chart_data = CategoryChartData()
    chart_data.categories = new_data[0]
    chart_data.add_series(chart.series[0].name ,tuple(series0),'0%')
    chart_data.add_series(chart.series[1].name ,tuple(series1),'0%')
    chart.replace_data(chart_data)
    for series in chart.plots[0].series:
        for i, val in enumerate(series.values):
            if val == 0:
                series.points[i].data_label.has_text_frame = True
        series.data_labels.show_value = True

def select_slides(presentation : Presentation, slides_to_select : list[int]):
    to_del = [i for i in range(len(presentation.slides)) if i not in slides_to_select]
    to_del.reverse()
    xml_slides = presentation.slides._sldIdLst  
    slides = list(xml_slides)
    for i in to_del:
        xml_slides.remove(slides[i]) 

def print_shapes_info(slide):
    for i, shape in enumerate(slide.shapes, 0):
        if shape.has_text_frame:
            print(f"{i}: text")
            print(f"\"{shape.text_frame.text}\"")
        if shape.has_table:
            print(f"{i}: table")
        if shape.has_chart:
            print(f"{i}: chart")

def create_slides_one_level(name: str, level: int, collector):   
    level_name = levels[level]
    level_table = levels_table[level]
    pres = Presentation("resources\\template.pptx")
    default_msg = "Недостаточно данных"
    select_slides(pres, [5])
    report = collector.get_person_report(name, level_name)

    #set positive/negative sides
    update_cards(pres.slides[0], report, levels[level])

    #set name
    update_text(pres.slides[0].shapes[1].text_frame, name)

    #set feedback

    feedback = default_msg if report[f"Обратная связь, {level_name.lower()}_count"] < 2 else "Обратная связь:\n" + report[f"Обратная связь, {level_name.lower()}_feedback"]
    update_text(pres.slides[0].shapes[3].text_frame, feedback)

    #set ratings
    update_table_cell(pres.slides[0].shapes[2].table, 0, 1, level_table)
    ratings = collector.get_person_ratings(name)
    rating = default_msg if ratings[level - 1][1] < 2 else f"{ratings[level - 1][0]:.1f}"
    update_table_cell(pres.slides[0].shapes[2].table, 1, 1, rating)
    overall_rating = collector.get_average_rating([4,11,20])["Overall"]
    update_table_cell(pres.slides[0].shapes[2].table, 1, 2, f"{overall_rating:.2f}")
    
    pres.save(f"outputs\\{name}.pptx")

def create_slides_two_levels(name: str, level: tuple[int,int],  collector):
    pres = Presentation("resources\\template.pptx")
    default_msg = "Недостаточно данных"
    select_slides(pres, [3,4])
    report = collector.get_person_report(name)

    #update charts
    update_chart(pres.slides[0].shapes[3].chart, collector.get_select_vals_for_plot(name, levels[level[0]]))
    update_text(pres.slides[0].shapes[3].chart.chart_title.text_frame, levels_chart[level[0]])
    update_chart(pres.slides[1].shapes[3].chart, collector.get_select_vals_for_plot(name, levels[level[1]]))
    update_text(pres.slides[1].shapes[3].chart.chart_title.text_frame, levels_chart[level[1]])

    #set name
    update_text(pres.slides[0].shapes[1].text_frame, name)
    update_text(pres.slides[1].shapes[1].text_frame, name)

    #set feedback
    feedback = default_msg if report[f"Обратная связь, {levels[level[0]].lower()}_count"] < 2 else (
        "Обратная связь:\n" + report[f"Обратная связь, {levels[level[0]].lower()}_feedback"])
    update_text(pres.slides[0].shapes[4].text_frame, feedback)

    feedback = default_msg if report[f"Обратная связь, {levels[level[1]].lower()}_count"] < 2 else (
        "Обратная связь:\n" + report[f"Обратная связь, {levels[level[1]].lower()}_feedback"])
    update_text(pres.slides[1].shapes[4].text_frame, feedback)

    #set ratings
    ratings = collector.get_person_ratings(name)

    def update_table(slide):
        update_table_cell(slide.shapes[2].table, 0, 1, levels_table[level[0]])
        rating = default_msg if ratings[level[0] - 1][1] < 2 else f"{ratings[level[0] - 1][0]:.1f}"
        update_table_cell(slide.shapes[2].table, 1, 1, rating)

        update_table_cell(slide.shapes[2].table, 0, 2, levels_table[level[1]])
        rating = default_msg if ratings[level[1] - 1][1] < 2 else f"{ratings[level[1] - 1][0]:.1f}"
        update_table_cell(slide.shapes[2].table, 1, 2, rating)

        overall_rating = collector.get_average_rating([4,11,20])["Overall"]
        update_table_cell(slide.shapes[2].table, 1, 3, f"{overall_rating:.2f}")
    
    update_table(pres.slides[0])
    update_table(pres.slides[1])
    
    pres.save(f"outputs\\{name}.pptx")

def create_slides_three_levels(name: str, collector):
    pres = Presentation("resources\\template.pptx")
    default_msg = "Недостаточно данных"
    select_slides(pres, [1,2])
    report = collector.get_person_report(name)

    #update charts
    update_chart(pres.slides[0].shapes[3].chart, collector.get_select_vals_for_plot(name, "Непосредственный руководитель"))
    update_chart(pres.slides[1].shapes[1].chart, collector.get_select_vals_for_plot(name, "Вышестоящий руководитель"))
    update_chart(pres.slides[1].shapes[2].chart, collector.get_select_vals_for_plot(name, "Функциональный руководитель"))

    #set name
    update_text(pres.slides[0].shapes[1].text_frame, name)
    update_text(pres.slides[1].shapes[5].text_frame, name)

    #set feedback
    feedback = default_msg if report["Обратная связь, непосредственный руководитель_count"] < 2 else "Обратная связь:\n" + report["Обратная связь, непосредственный руководитель_feedback"]
    update_text(pres.slides[0].shapes[4].text_frame, feedback)

    feedback = default_msg if report["Обратная связь, вышестоящий руководитель_count"] < 2 else "Обратная связь:\n" + report["Обратная связь, вышестоящий руководитель_feedback"]
    update_text(pres.slides[1].shapes[3].text_frame, feedback)

    feedback = default_msg if report["Обратная связь, функциональный руководитель_count"] < 2 else "Обратная связь:\n" + report["Обратная связь, функциональный руководитель_feedback"]
    update_text(pres.slides[1].shapes[4].text_frame, feedback)

    #set ratings
    ratings = collector.get_person_ratings(name)
    rating = default_msg if ratings[0][1] < 2 else f"{ratings[0][0]:.1f}"
    update_table_cell(pres.slides[0].shapes[2].table, 1, 1, rating)

    rating = default_msg if ratings[1][1] < 2 else f"{ratings[1][0]:.1f}"
    update_table_cell(pres.slides[0].shapes[2].table, 1, 2, rating)

    rating = default_msg if ratings[2][1] < 2 else f"{ratings[2][0]:.1f}"
    update_table_cell(pres.slides[0].shapes[2].table, 1, 3, rating)

    overall_rating = collector.get_average_rating([4,11,20])["Overall"]
    update_table_cell(pres.slides[0].shapes[2].table, 1, 4, f"{overall_rating:.2f}")
    
    pres.save(f"outputs\\{name}.pptx")
